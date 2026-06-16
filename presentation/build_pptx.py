#!/usr/bin/env python3
"""
Build a professional, academic PowerPoint (.pptx) presentation for the thesis:

    "Making AI More Trustworthy: Context Labelling and Original-Source
     Tracking for Clinical Question Answering"  (GRAPES-SHAP framework)

The deck embeds the publication figures from outputs/figures/paper and is
fully self-contained (images are copied into the .pptx).  The output also
uploads cleanly to Google Slides (File > Import / upload).

Run:
    python presentation/build_pptx.py
Output:
    presentation/GRAPES_SHAP_Thesis_Presentation.pptx
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ───────────────────────── paths ─────────────────────────
ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "presentation" / "assets" / "figures"
if not FIG.exists():
    FIG = ROOT / "outputs" / "figures" / "paper"
OUT = ROOT / "presentation" / "GRAPES_SHAP_Thesis_Presentation.pptx"


def f(name: str) -> str:
    """Resolve a figure path by stem (png preferred)."""
    p = FIG / f"{name}.png"
    return str(p)


# ──────────────────── academic colour palette ────────────────────
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)   # primary
STEEL  = RGBColor(0x5B, 0x7B, 0x96)   # secondary
TEAL   = RGBColor(0x2E, 0x6E, 0x6A)   # proposed / positive
SIENNA = RGBColor(0x9C, 0x5A, 0x33)   # warm accent
PLUM   = RGBColor(0x6E, 0x5A, 0x74)   # reasoning
GRAY   = RGBColor(0x5F, 0x6B, 0x76)   # neutral
GOLD   = RGBColor(0xB0, 0x89, 0x2F)   # highlight
INK    = RGBColor(0x21, 0x29, 0x33)   # body text
FILL   = RGBColor(0xEA, 0xEE, 0xF2)   # light fill
LIGHT  = RGBColor(0xF6, 0xF8, 0xFA)   # slide background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MIST   = RGBColor(0xDD, 0xE5, 0xEC)   # rule / subtle line
DEEP   = RGBColor(0x15, 0x28, 0x42)   # darker navy for dividers

HEAD_FONT = "Georgia"
BODY_FONT = "Segoe UI"

# ───────────────────────── deck setup ─────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

_page = {"n": 0}


# ───────────────────────── helpers ─────────────────────────
def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def new_slide(bg=LIGHT):
    slide = prs.slides.add_slide(BLANK)
    b = slide.background
    b.fill.solid()
    b.fill.fore_color.rgb = bg
    # fade transition
    try:
        slide._element.append(parse_xml(
            '<p:transition xmlns:p="http://schemas.openxmlformats.org/'
            'presentationml/2006/main" spd="med"><p:fade/></p:transition>'))
    except Exception:
        pass
    return slide


def rect(slide, x, y, w, h, color, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
         radius=None):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    _no_shadow(sh)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def _apply(run, props):
    fp = props.get
    run.font.name = fp("name", BODY_FONT)
    run.font.size = Pt(fp("size", 16))
    run.font.bold = fp("bold", False)
    run.font.italic = fp("italic", False)
    run.font.color.rgb = fp("color", INK)
    sp = fp("spacing")
    if sp is not None:
        run.font._rPr.set("spc", str(int(sp)))


def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts {align, space_after, space_before, line, level,
    runs:[(text, props)]}."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if para.get("space_after") is not None:
            p.space_after = Pt(para["space_after"])
        if para.get("space_before") is not None:
            p.space_before = Pt(para["space_before"])
        if para.get("line") is not None:
            p.line_spacing = para["line"]
        if para.get("level"):
            p.level = para["level"]
        for txt, props in para["runs"]:
            r = p.add_run()
            r.text = txt
            _apply(r, props)
    return tb


def footer(slide, accent=TEAL, dark=False):
    _page["n"] += 1
    col = RGBColor(0xC7, 0xD2, 0xDD) if dark else GRAY
    rect(slide, 0.55, 7.06, 12.23, 0.012, MIST if not dark else RGBColor(0x33, 0x49, 0x66))
    textbox(slide, 0.55, 7.12, 9.0, 0.3,
            [{"runs": [("Making AI More Trustworthy  ·  Clinical QA via Context "
                        "Labelling & Source Tracking", dict(size=9, color=col))]}])
    textbox(slide, 11.0, 7.12, 1.78, 0.3,
            [{"align": PP_ALIGN.RIGHT,
              "runs": [(f"{_page['n']:02d}", dict(size=9, color=accent, bold=True))]}])


def header(slide, kicker, title, accent=TEAL):
    # left accent bar
    rect(slide, 0.0, 0.0, 0.22, 7.5, accent)
    # kicker
    textbox(slide, 0.6, 0.42, 11.5, 0.35,
            [{"runs": [(kicker.upper(),
                        dict(size=12.5, color=accent, bold=True, spacing=240,
                             name=BODY_FONT))]}])
    # title
    textbox(slide, 0.58, 0.74, 12.2, 0.95,
            [{"runs": [(title, dict(size=29, color=NAVY, bold=True,
                                    name=HEAD_FONT))]}])
    # underline accent
    rect(slide, 0.62, 1.62, 1.7, 0.06, accent)
    footer(slide, accent)


def place_image(slide, path, x, y, w, h, caption=None, frame=True,
                caption_color=GRAY):
    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 1600, 1000
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        dw, dh = w, w / ar
    else:
        dh, dw = h, h * ar
    px = x + (w - dw) / 2
    py = y + (h - dh) / 2
    if frame:
        rect(slide, px - 0.06, py - 0.06, dw + 0.12, dh + 0.12, WHITE,
             line=MIST, lw=1.0)
    slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(dw), Inches(dh))
    if caption:
        textbox(slide, x, y + h + 0.02, w, 0.32,
                [{"align": PP_ALIGN.CENTER,
                  "runs": [(caption, dict(size=10.5, italic=True,
                                          color=caption_color))]}])
    return px, py, dw, dh


def bullets(slide, x, y, w, h, items, size=16, gap=9, accent=TEAL,
            anchor=MSO_ANCHOR.TOP):
    """items: list of (lead, body, level). lead bold (may be None)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, it in enumerate(items):
        lead, body, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.06
        p.level = lvl
        marker = p.add_run()
        marker.text = ("▸  " if lvl == 0 else "–  ")
        _apply(marker, dict(size=size, bold=True,
                            color=accent if lvl == 0 else GRAY))
        if lead:
            r = p.add_run()
            r.text = lead + (" " if body else "")
            _apply(r, dict(size=size, bold=True, color=NAVY))
        if body:
            r2 = p.add_run()
            r2.text = body
            _apply(r2, dict(size=size, color=INK))
    return tb


def cards(slide, y, items, x0=0.6, total_w=12.13, h=1.55, gap=0.28, num_size=33,
          lab_size=11.5):
    """items: list of (number, label, color). Layout adapts to card height."""
    n = len(items)
    w = (total_w - gap * (n - 1)) / n
    x = x0
    num_h = h * 0.56
    for num, label, color in items:
        rect(slide, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.08)
        textbox(slide, x + 0.1, y + 0.06, w - 0.2, num_h,
                [{"align": PP_ALIGN.CENTER,
                  "runs": [(num, dict(size=num_size, bold=True, color=WHITE,
                                      name=HEAD_FONT))]}],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, x + 0.12, y + num_h + 0.02, w - 0.24, h - num_h - 0.04,
                [{"align": PP_ALIGN.CENTER,
                  "runs": [(label, dict(size=lab_size, color=WHITE))]}],
                anchor=MSO_ANCHOR.TOP)
        x += w + gap


def table(slide, x, y, w, h, data, col_w=None, size=12.5, header=True,
          header_color=NAVY, first_col_bold=False, align_first_left=True):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    tbl = gt.table
    # strip default banded style for a clean academic look
    tbl.first_row = header
    tbl.horz_banding = False
    if col_w:
        for c, cw in enumerate(col_w):
            tbl.columns[c].width = Inches(cw)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = str(data[r][c])
            run = para.runs[0] if para.runs else para.add_run()
            run.font.name = BODY_FONT
            run.font.size = Pt(size)
            if c == 0 and align_first_left:
                para.alignment = PP_ALIGN.LEFT
            else:
                para.alignment = PP_ALIGN.CENTER
            if r == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                run.font.color.rgb = WHITE
                run.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if (r % 2 == 1) else FILL
                run.font.color.rgb = INK
                if (c == 0 and first_col_bold):
                    run.font.bold = True
                    run.font.color.rgb = NAVY
    return tbl


def chip(slide, x, y, w, text_, color, txt_color=WHITE, size=11.5, h=0.42):
    rect(slide, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    textbox(slide, x, y + 0.02, w, h - 0.02,
            [{"align": PP_ALIGN.CENTER,
              "runs": [(text_, dict(size=size, bold=True, color=txt_color))]}],
            anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════ 1 · TITLE ════════════════════════
def slide_title():
    s = new_slide(DEEP)
    # background bands
    rect(s, 0, 0, SW, SH, DEEP)
    rect(s, 0, 0, SW, 0.16, TEAL)
    rect(s, 0, SH - 0.16, SW, 0.16, GOLD)
    # decorative side accents
    rect(s, 0, 1.0, 0.16, 5.5, TEAL)
    # kicker
    textbox(s, 1.0, 1.18, 11.3, 0.4,
            [{"runs": [("M.SC. / THESIS DEFENCE  ·  TRUSTWORTHY CLINICAL AI",
                        dict(size=13, color=RGBColor(0x9F, 0xC4, 0xC2), bold=True,
                             spacing=260))]}])
    # title
    textbox(s, 0.95, 1.85, 11.5, 2.3,
            [{"line": 1.04,
              "runs": [("Making AI More Trustworthy", dict(size=46, bold=True,
                        color=WHITE, name=HEAD_FONT))]},
             {"line": 1.04, "space_before": 4,
              "runs": [("Context Labelling & Original-Source Tracking\n"
                        "for Clinical Question Answering",
                        dict(size=27, color=RGBColor(0xCF, 0xDD, 0xEA),
                             name=HEAD_FONT))]}])
    # subtitle line
    rect(s, 1.0, 4.5, 3.2, 0.05, GOLD)
    textbox(s, 1.0, 4.66, 11.0, 0.6,
            [{"runs": [("A World-Model–Guided, Calibrated & Interpretable "
                        "Retrieval-Augmented Generation Framework",
                        dict(size=15.5, italic=True,
                             color=RGBColor(0xB9, 0xCC, 0xDB)))]}])
    # meta chips
    chip(s, 1.0, 5.5, 2.9, "GRAPES-SHAP", TEAL, size=12.5)
    chip(s, 4.05, 5.5, 3.0, "10.1 M Parameters", STEEL, size=12.5)
    chip(s, 7.2, 5.5, 3.0, "100,000 Samples", PLUM, size=12.5)
    chip(s, 10.35, 5.5, 1.95, "DeepSeek LLM", SIENNA, size=11)
    # author block
    textbox(s, 1.0, 6.35, 11.0, 0.7,
            [{"runs": [("Author Name", dict(size=15, bold=True, color=WHITE))]},
             {"runs": [("Department / Institution  ·  Supervisor: [Name]  ·  2026",
                        dict(size=12, color=RGBColor(0x9C, 0xB1, 0xC4)))]}])


# ════════════════════════ DIVIDER ════════════════════════
def divider(num, title, subtitle, accent=TEAL):
    s = new_slide(DEEP)
    rect(s, 0, 0, SW, SH, DEEP)
    rect(s, 0, 0, 0.22, SH, accent)
    rect(s, 0.9, 2.5, 2.2, 0.07, accent)
    textbox(s, 0.85, 1.7, 6.0, 1.4,
            [{"runs": [(num, dict(size=104, bold=True,
                        color=RGBColor(0x2A, 0x44, 0x68), name=HEAD_FONT))]}])
    textbox(s, 0.9, 2.75, 11.0, 1.2,
            [{"runs": [(title, dict(size=36, bold=True, color=WHITE,
                                    name=HEAD_FONT))]}])
    textbox(s, 0.92, 3.95, 10.8, 0.9,
            [{"runs": [(subtitle, dict(size=16, italic=True,
                        color=RGBColor(0xB9, 0xCC, 0xDB)))]}])
    # page tick
    _page["n"] += 1


# ════════════════════════ 2 · OUTLINE ════════════════════════
def slide_outline():
    s = new_slide()
    header(s, "Presentation Roadmap", "Outline", NAVY)
    items = [
        ("01", "Motivation & Problem", "Why fluent LLMs are not yet trustworthy", TEAL),
        ("02", "Objectives & Contributions", "Six goals; two core trust mechanisms", STEEL),
        ("03", "Approach & Architecture", "Pipeline, world model, neural core", PLUM),
        ("04", "Data & Preprocessing", "DDXPlus · MedMCQA · MedQA (100k)", SIENNA),
        ("05", "Training & Evaluation", "Convergence, calibration, generalisation", TEAL),
        ("06", "Results & Impact", "Head-to-head vs strong RAG baseline", NAVY),
        ("07", "Planning · Risk · Economics", "Engineering & responsible-AI view", GRAY),
        ("08", "Conclusion & Future Work", "Findings, limits, next steps", GOLD),
    ]
    x0, y0 = 0.7, 2.1
    cw, ch = 6.0, 1.12
    gapx, gapy = 0.25, 0.18
    for i, (n, t, sub, col) in enumerate(items):
        col_i = i % 2
        row_i = i // 2
        x = x0 + col_i * (cw + gapx)
        y = y0 + row_i * (ch + gapy)
        rect(s, x, y, cw, ch, WHITE, line=MIST, lw=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
        rect(s, x, y, 0.12, ch, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        textbox(s, x + 0.28, y + 0.12, 1.1, ch - 0.2,
                [{"runs": [(n, dict(size=26, bold=True, color=col,
                                    name=HEAD_FONT))]}],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + 1.35, y + 0.16, cw - 1.5, 0.5,
                [{"runs": [(t, dict(size=15.5, bold=True, color=NAVY))]}])
        textbox(s, x + 1.35, y + 0.62, cw - 1.5, 0.4,
                [{"runs": [(sub, dict(size=11.5, color=GRAY))]}])


# ════════════════════════ 3 · PROBLEM ════════════════════════
def slide_problem():
    s = new_slide()
    header(s, "Section 01 · Motivation", "The Trust Gap in Clinical LLMs", TEAL)
    bullets(s, 0.7, 2.0, 6.2, 4.6, [
        ("Fluency ≠ trustworthiness.", "LLMs encode rich medical knowledge "
         "yet remain hard to deploy safely.", 0),
        ("Hallucination.", "Confident, well-formed claims unsupported by any "
         "source.", 0),
        ("Opaque confidence.", "A 0.95-confidence answer and a borderline one "
         "look identical.", 0),
        ("No provenance.", "Clinicians cannot see which evidence actually drove "
         "the answer.", 0),
        ("RAG helps, but is reactive.", "Retrieve–concatenate–generate delegates "
         "all reasoning & self-checking to one opaque pass.", 0),
    ], size=16.5, gap=12)
    # right call-out panel
    rect(s, 7.3, 2.0, 5.4, 4.5, FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    rect(s, 7.3, 2.0, 5.4, 0.7, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    textbox(s, 7.5, 2.12, 5.0, 0.5,
            [{"runs": [("Three properties canonical RAG lacks",
                        dict(size=14.5, bold=True, color=WHITE))]}],
            anchor=MSO_ANCHOR.MIDDLE)
    panel = [
        ("Prospective reasoning", "Simulate how a patient responds before "
         "committing to an action.", TEAL),
        ("Calibrated self-knowledge", "Know what it does not know; expose "
         "honest uncertainty.", STEEL),
        ("Provenance", "State which documents were decisive — auditable, "
         "not decorative.", SIENNA),
    ]
    yy = 2.95
    for t, b, c in panel:
        rect(s, 7.55, yy, 0.12, 1.0, c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        textbox(s, 7.85, yy + 0.02, 4.7, 0.4,
                [{"runs": [(t, dict(size=14, bold=True, color=NAVY))]}])
        textbox(s, 7.85, yy + 0.42, 4.7, 0.6,
                [{"runs": [(b, dict(size=12, color=INK))]}])
        yy += 1.15


# ════════════════════════ 4 · OBJECTIVES ════════════════════════
def slide_objectives():
    s = new_slide()
    header(s, "Section 01 · Goals", "Research Objectives", STEEL)
    obj = [
        ("Augment retrieval with a world model", "action-conditioned latent "
         "dynamics for prospective, tree-of-thought planning", TEAL),
        ("Calibrated uncertainty", "deep ensemble with epistemic / aleatoric "
         "decomposition", STEEL),
        ("Original-source tracking", "Shapley-value evidence attribution with "
         "explicit citations", SIENNA),
        ("Compact & efficient core", "~10 M parameters, < 20 min training on one "
         "consumer GPU", PLUM),
        ("Rigorous trust evaluation", "calibration, generalisation, significance, "
         "controlled baseline comparison", NAVY),
        ("Honest scope & limitations", "explicit boundaries of validity for "
         "responsible interpretation", GRAY),
    ]
    x0, y0 = 0.7, 2.05
    cw, ch = 5.95, 1.32
    gx, gy = 0.3, 0.22
    for i, (t, b, c) in enumerate(obj):
        x = x0 + (i % 2) * (cw + gx)
        y = y0 + (i // 2) * (ch + gy)
        rect(s, x, y, cw, ch, WHITE, line=MIST, lw=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        rect(s, x, y, cw, 0.1, c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        # number badge
        rect(s, x + 0.28, y + 0.34, 0.62, 0.62, c,
             shape=MSO_SHAPE.OVAL)
        textbox(s, x + 0.28, y + 0.36, 0.62, 0.6,
                [{"align": PP_ALIGN.CENTER,
                  "runs": [(f"{i+1}", dict(size=18, bold=True, color=WHITE,
                                           name=HEAD_FONT))]}],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + 1.1, y + 0.22, cw - 1.3, 0.5,
                [{"runs": [(t, dict(size=14.5, bold=True, color=NAVY))]}])
        textbox(s, x + 1.1, y + 0.66, cw - 1.3, 0.6,
                [{"runs": [(b, dict(size=11.8, color=INK))]}])


# ════════════════════════ 5 · TWO MECHANISMS ════════════════════════
def slide_mechanisms():
    s = new_slide()
    header(s, "Section 02 · Core Idea", "Two Coupled Trust Mechanisms", PLUM)
    # left card
    rect(s, 0.7, 2.05, 5.85, 4.5, WHITE, line=MIST, lw=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    rect(s, 0.7, 2.05, 5.85, 0.95, TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    textbox(s, 0.95, 2.18, 5.4, 0.7,
            [{"runs": [("①  Context Labelling", dict(size=20, bold=True,
                        color=WHITE, name=HEAD_FONT))]}],
            anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 1.0, 3.2, 5.3, 3.2, [
        (None, "Attach a structured, machine-checkable record to every answer:", 0),
        (None, "Retrieved evidence set", 1),
        (None, "Causal knowledge-graph relations", 1),
        (None, "Predicted diagnostic state", 1),
        (None, "Calibrated uncertainty estimate", 1),
    ], size=14.5, gap=9, accent=TEAL)
    # right card
    rect(s, 6.8, 2.05, 5.85, 4.5, WHITE, line=MIST, lw=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    rect(s, 6.8, 2.05, 5.85, 0.95, SIENNA, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    textbox(s, 7.05, 2.18, 5.4, 0.7,
            [{"runs": [("②  Original-Source Tracking", dict(size=20, bold=True,
                        color=WHITE, name=HEAD_FONT))]}],
            anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 7.1, 3.2, 5.3, 3.2, [
        (None, "Trace each claim back to the documents that supported it:", 0),
        (None, "Shapley value per retrieved document", 1),
        (None, "Quantifies marginal contribution", 1),
        (None, "Explicit, auditable citations", 1),
        (None, "Converts an opaque generator into an accountable assistant", 1),
    ], size=14.5, gap=9, accent=SIENNA)


# ════════════════════════ 6 · SYSTEM ARCHITECTURE ════════════════════════
def slide_fig(kicker, title, fig_name, caption, accent=NAVY, note=None,
              img_box=(0.7, 1.95, 11.93, 4.55)):
    s = new_slide()
    header(s, kicker, title, accent)
    x, y, w, h = img_box
    place_image(s, f(fig_name), x, y, w, h, caption=caption)
    if note:
        textbox(s, 0.7, 6.62, 11.9, 0.4,
                [{"align": PP_ALIGN.CENTER,
                  "runs": [(note, dict(size=12, color=GRAY, italic=True))]}])
    return s


# ════════════════════════ DATASET TABLE ════════════════════════
def slide_datasets():
    s = new_slide()
    header(s, "Section 03 · Data", "Datasets & Corpus", SIENNA)
    data = [
        ["Dataset", "Role", "Samples", "Source"],
        ["DDXPlus", "Diagnosis trajectories (train/val/test)", "80k / 10k / 10k",
         "HuggingFace"],
        ["MedMCQA", "Retrieval corpus (evidence)", "50,000 docs", "HuggingFace"],
        ["MedQA (USMLE)", "Evaluation queries", "1,000", "HuggingFace"],
        ["Total", "Integrated multi-source corpus", "100,000", "Open / free"],
    ]
    table(s, 0.7, 2.05, 8.2, 2.4, data,
          col_w=[2.0, 3.5, 1.7, 1.0], size=13.5, first_col_bold=True)
    # right facts
    rect(s, 9.15, 2.05, 3.5, 4.45, FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    textbox(s, 9.4, 2.2, 3.1, 0.4,
            [{"runs": [("DDXPlus at a glance", dict(size=14, bold=True,
                        color=NAVY))]}])
    facts = [
        ("49", "pathology classes (long-tailed)"),
        ("223", "normalised evidence codes (from 504)"),
        ("39.6 yr", "mean patient age"),
        ("19", "median evidences per patient"),
        ("100%", "records retained after cleaning"),
    ]
    yy = 2.7
    for v, lab in facts:
        textbox(s, 9.4, yy, 3.05, 0.4,
                [{"runs": [(v + "  ", dict(size=17, bold=True, color=SIENNA,
                            name=HEAD_FONT)),
                           (lab, dict(size=11, color=INK))]}])
        yy += 0.74
    bullets(s, 0.7, 4.7, 8.2, 1.9, [
        ("Zero-barrier.", "All datasets are public, free, and auto-downloaded "
         "from HuggingFace — no credentials required.", 0),
        ("Integration.", "Heterogeneous sources joined into a single 64-dim "
         "observation + 5-dim outcome tensor format.", 0),
    ], size=14, gap=10, accent=SIENNA)


# ════════════════════════ MODEL CONFIG ════════════════════════
def slide_config():
    s = new_slide()
    header(s, "Section 03 · Configuration", "Model & Training Configuration", PLUM)
    arch = [
        ["Architecture", "Value"],
        ["Observation / action dim", "64 / 50"],
        ["Latent / hidden dim", "256 / 512"],
        ["Graph nodes / node dim", "20 / 128"],
        ["Transformer layers / heads", "3 / 8"],
        ["GRU layers", "3"],
        ["Ensemble members", "5"],
        ["Sequence length", "8"],
    ]
    train = [
        ["Training", "Value"],
        ["WM / predictor epochs", "15 / 10"],
        ["WM / predictor LR", "2e-4 / 1e-3"],
        ["Batch size", "64"],
        ["Optimiser / AMP", "Adam / fp16"],
        ["Top-k retrieval", "6"],
        ["SHAP permutations", "32"],
        ["Plan horizon / beam", "4 / 8"],
    ]
    table(s, 0.7, 2.05, 5.7, 3.95, arch, col_w=[3.7, 2.0], size=13,
          first_col_bold=True, header_color=NAVY)
    table(s, 6.65, 2.05, 5.95, 3.95, train, col_w=[3.8, 2.15], size=13,
          first_col_bold=True, header_color=TEAL)
    cards(s, 6.1, [
        ("10.13 M", "Trainable parameters", NAVY),
        ("18.8 min", "Training time", TEAL),
        ("RTX 4080S", "Single GPU", STEEL),
        ("PyTorch", "2.6.0 + CUDA", GRAY),
    ], h=0.86, num_size=21, lab_size=10.5)


# ════════════════════════ KEY RESULTS (impact) ════════════════════════
def slide_key_results():
    s = new_slide()
    header(s, "Section 04 · Headline", "Key Results & Impact", NAVY)
    cards(s, 2.0, [
        ("0.039", "World-model MAE", TEAL),
        ("0.029", "Calibration error (ECE)", STEEL),
        ("0.80", "1σ coverage", PLUM),
        ("< 0.002", "Train→test MAE gap", SIENNA),
    ], h=1.5, num_size=34)
    cards(s, 3.85, [
        ("0.70 → 0.97", "Clinical concept coverage", NAVY),
        ("0.50 → 0.84", "Answer completeness", TEAL),
        ("2.0 → 5.1", "Evidence citations / answer", STEEL),
        ("9 / 10", "Vignettes won vs baseline", GOLD),
    ], h=1.5, num_size=26)
    # bottom strip
    rect(s, 0.7, 5.7, 11.93, 1.0, FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    textbox(s, 0.95, 5.83, 11.5, 0.8,
            [{"runs": [("Statistically significant ", dict(size=14.5, bold=True,
                        color=NAVY)),
                       ("improvement over a strong hybrid-RAG baseline — Wilcoxon "
                        "signed-rank p = 0.046, mean +0.27 (95% CI ±0.21) — while "
                        "adding calibrated uncertainty and per-document source "
                        "tracking the baseline cannot provide.",
                        dict(size=14.5, color=INK))]}],
            anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════ EVAL METRICS table ════════════════════════
def slide_eval_metrics():
    s = new_slide()
    header(s, "Section 04 · Evaluation", "Held-out Evaluation Metrics", TEAL)
    data = [
        ["Metric", "Value", "Target", "Status"],
        ["MAE (outcome regression)", "0.039", "lower better", "✓"],
        ["RMSE", "0.075", "lower better", "✓"],
        ["1σ coverage (calibration)", "0.80", "≈ 0.68–0.95", "✓"],
        ["Expected Calibration Error", "0.029", "< 0.05", "✓"],
        ["Diagnosis accuracy (top-k)", "0.755", "—", "✓"],
        ["F1-macro", "0.172", "long-tail (49 cls)", "△"],
        ["Mean |SHAP|", "0.744", "—", "—"],
    ]
    table(s, 0.7, 2.05, 7.0, 4.3, data, col_w=[3.3, 1.3, 1.7, 0.7], size=13,
          first_col_bold=True)
    place_image(s, f("fig7_metrics_dashboard"), 7.95, 2.05, 4.7, 3.7)
    textbox(s, 0.7, 6.5, 11.9, 0.45,
            [{"align": PP_ALIGN.CENTER,
              "runs": [("Headline: 1σ coverage 0.80 with ECE 0.029 — the "
                        "ensemble's uncertainty is trustworthy, not "
                        "over-confident.",
                        dict(size=12.5, italic=True, color=GRAY))]}])


# ════════════════════════ CONTRIBUTIONS ════════════════════════
def slide_contributions():
    s = new_slide()
    header(s, "Section 08 · Closing", "Contributions to the Field", NAVY)
    items = [
        ("Trust-oriented reframing of RAG", "organised around context labelling "
         "and original-source tracking.", TEAL),
        ("Compact world-model core", "supervised reward shaping + beam-search "
         "planning adds prospective reasoning to retrieval.", STEEL),
        ("Native calibrated uncertainty", "explicit epistemic / aleatoric "
         "decomposition via a deep ensemble.", PLUM),
        ("Shapley source tracking", "attributes answers to individual documents "
         "for auditability.", SIENNA),
        ("Honest, reproducible science", "all code, configs & figure scripts "
         "released; scope stated explicitly.", GRAY),
    ]
    y = 2.05
    for i, (t, b, c) in enumerate(items):
        rect(s, 0.7, y, 11.95, 0.82, WHITE, line=MIST, lw=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
        rect(s, 0.7, y, 0.12, 0.82, c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        rect(s, 1.0, y + 0.16, 0.5, 0.5, c, shape=MSO_SHAPE.OVAL)
        textbox(s, 1.0, y + 0.17, 0.5, 0.5,
                [{"align": PP_ALIGN.CENTER,
                  "runs": [(f"{i+1}", dict(size=15, bold=True, color=WHITE,
                            name=HEAD_FONT))]}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, 1.7, y + 0.13, 10.7, 0.6,
                [{"runs": [(t + " — ", dict(size=14.5, bold=True, color=NAVY)),
                           (b, dict(size=13.5, color=INK))]}],
                anchor=MSO_ANCHOR.MIDDLE)
        y += 0.92


# ════════════════════════ LIMITATIONS & FUTURE ════════════════════════
def slide_future():
    s = new_slide()
    header(s, "Section 08 · Outlook", "Limitations & Future Work", SIENNA)
    rect(s, 0.7, 2.05, 5.85, 4.5, FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    rect(s, 0.7, 2.05, 5.85, 0.7, GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    textbox(s, 0.95, 2.16, 5.4, 0.5,
            [{"runs": [("Honest Limitations", dict(size=17, bold=True,
                        color=WHITE, name=HEAD_FONT))]}], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 1.0, 2.95, 5.3, 3.5, [
        ("Semi-synthetic dynamics.", "trajectories reveal evidence sequentially; "
         "transitions near-deterministic.", 0),
        ("Evidence-acquisition actions", "rather than physical interventions.", 0),
        ("Stochastic causal prior", "rather than a curated ontology.", 0),
        ("Low macro-F1", "from fine-grained, imbalanced 49-class ranking.", 0),
    ], size=13.5, gap=10, accent=GRAY)
    rect(s, 6.8, 2.05, 5.85, 4.5, FILL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    rect(s, 6.8, 2.05, 5.85, 0.7, TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    textbox(s, 7.05, 2.16, 5.4, 0.5,
            [{"runs": [("Future Directions", dict(size=17, bold=True,
                        color=WHITE, name=HEAD_FONT))]}], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 7.1, 2.95, 5.3, 3.5, [
        ("Logged intervention data", "(e.g. ICU records) for true treatment "
         "planning.", 0),
        ("Curated ontology", "(e.g. SNOMED-CT) to strengthen graph reasoning.", 0),
        ("Class re-weighting / focal loss", "to raise macro-F1 without hurting "
         "calibration.", 0),
        ("Clinician-in-the-loop study", "to test reduction of automation bias.", 0),
    ], size=13.5, gap=10, accent=TEAL)


# ════════════════════════ CONCLUSION ════════════════════════
def slide_conclusion():
    s = new_slide(DEEP)
    rect(s, 0, 0, SW, SH, DEEP)
    rect(s, 0, 0, SW, 0.16, TEAL)
    rect(s, 0, SH - 0.16, SW, 0.16, GOLD)
    rect(s, 0.9, 1.0, 2.2, 0.07, TEAL)
    textbox(s, 0.9, 1.15, 11.0, 0.9,
            [{"runs": [("Conclusion", dict(size=40, bold=True, color=WHITE,
                        name=HEAD_FONT))]}])
    bullets(s, 0.95, 2.3, 11.4, 3.2, [
        (None, "A compact 10.1 M-parameter framework makes clinical QA more "
         "trustworthy via context labelling and original-source tracking.", 0),
        (None, "Accurate (MAE 0.039), well-calibrated (ECE 0.029) and "
         "generalising (train→test gap < 0.002 MAE).", 0),
        (None, "Significantly better coverage, completeness and evidence "
         "grounding than a strong RAG baseline (p = 0.046).", 0),
        (None, "Adds calibrated uncertainty and per-document attribution that "
         "the baseline cannot provide — at negligible cost.", 0),
    ], size=16.5, gap=13, accent=TEAL)
    # closing line
    rect(s, 0.95, 5.85, 11.4, 0.9, RGBColor(0x21, 0x3A, 0x5C),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    textbox(s, 1.2, 5.97, 11.0, 0.7,
            [{"runs": [("Graph-grounded retrieval + world-model planning + "
                        "ensemble uncertainty ⇒ more accurate, better-grounded, "
                        "and more trustworthy clinical recommendations.",
                        dict(size=14.5, italic=True, color=WHITE))]}],
            anchor=MSO_ANCHOR.MIDDLE)
    _page["n"] += 1


# ════════════════════════ THANK YOU ════════════════════════
def slide_thanks():
    s = new_slide(DEEP)
    rect(s, 0, 0, SW, SH, DEEP)
    rect(s, 0, 1.0, 0.16, 5.5, TEAL)
    textbox(s, 1.0, 2.3, 11.0, 1.4,
            [{"runs": [("Thank You", dict(size=58, bold=True, color=WHITE,
                        name=HEAD_FONT))]}])
    rect(s, 1.05, 3.65, 3.0, 0.06, GOLD)
    textbox(s, 1.05, 3.9, 11.0, 0.7,
            [{"runs": [("Questions & Discussion", dict(size=22, italic=True,
                        color=RGBColor(0xCF, 0xDD, 0xEA)))]}])
    chip(s, 1.05, 4.95, 3.4, "GRAPES-SHAP Framework", TEAL, size=12.5)
    chip(s, 4.6, 4.95, 3.2, "Reproducible · Open", STEEL, size=12.5)
    chip(s, 7.95, 4.95, 3.6, "Trustworthy Clinical AI", SIENNA, size=12.5)
    textbox(s, 1.05, 6.0, 11.0, 0.5,
            [{"runs": [("Author Name  ·  Department / Institution  ·  2026",
                        dict(size=13, color=RGBColor(0x9C, 0xB1, 0xC4)))]}])
    _page["n"] += 1


# ════════════════════════ MULTI-METHOD RAG COMPARISON ════════════════════════
def slide_multirag():
    s = new_slide()
    header(s, "Section 04 · Advanced RAG",
           "Benchmarked Against Four Advanced-RAG Pipelines", NAVY)
    place_image(s, f("fig24_multi_method"), 1.25, 1.78, 10.8, 3.45,
                caption="Live run — identical MedMCQA corpus, same DeepSeek LLM "
                        "and rubric across all five systems; only the "
                        "retrieval / reasoning stack varies.")
    bullets(s, 0.7, 5.55, 12.13, 1.35, [
        ("Stronger retrieval alone is not enough.",
         "HyDE leads retrieval-only coverage (0.87); hybrid fusion and "
         "cross-encoder + MMR do not uniformly help, and all plateau near 0.50 "
         "on answer structure.", 0),
        ("Our reasoning stack drives the trust-relevant gains.",
         "Coverage 0.97 (+10 pts over HyDE), structure 0.84, grounding 5.1 "
         "citations — plus calibrated uncertainty, SHAP attribution and plan "
         "simulation that no baseline provides.", 0),
    ], size=13, gap=8, accent=NAVY)
    return s


# ════════════════════════ WORKED ANSWER DEMO ════════════════════════
def slide_demo():
    s = new_slide()
    header(s, "Section 04 · Answer Demo",
           "Worked Answer Comparisons — Baseline vs Ours", NAVY)
    data = [
        ["Scenario", "Baseline RAG", "GRAPES-SHAP (ours)"],
        ["P1 · STEMI (chest pain, ST-elevation II/III/aVF)",
         "Inferior-wall STEMI; reperfusion — PCI or thrombolysis.   "
         "[coverage 0.67 · 2 cites · conf 0.95]",
         "STEMI confirmed; PCI vs thrombolysis weighed (PCI: lower bleed risk); "
         "full supportive bundle; plan score -0.29.   "
         "[coverage 1.00 · 4 cites · conf 0.70]"],
        ["P2 · Septic shock (fever, lactate 3.2, BP 88/52)",
         "Septic shock; fluids, oxygen, norepinephrine if fluids fail.   "
         "[coverage 0.33 · 1 cite · conf 0.90]",
         "Three pillars — source control + early broad-spectrum antibiotics + "
         "pressors; reconciles dopamine-vs-norepinephrine evidence.   "
         "[coverage 1.00 · 6 cites · conf 0.70]"],
        ["P4 · Ischaemic stroke (AF, NIHSS 14, 1.5 h onset)",
         "Ischaemic stroke (left MCA); IV alteplase within 4.5 h.   "
         "[coverage 0.33 · 2 cites · conf 0.95]",
         "Localisation → cardioembolic AF → eligibility → tPA; SHAP ranks the "
         "time-window & no-haemorrhage evidence as decisive.   "
         "[coverage 1.00 · 4 cites · conf 0.70]"],
    ]
    table(s, 0.55, 1.95, 12.25, 4.2, data,
          col_w=[2.95, 4.3, 5.0], size=10.5, first_col_bold=True)
    textbox(s, 0.7, 6.4, 12.13, 0.5,
            [{"align": PP_ALIGN.CENTER,
              "runs": [("Same top-line diagnosis — but ours surfaces explicit "
                        "trade-offs, the outcome-critical missing step, calibrated "
                        "confidence and per-evidence attribution.",
                        dict(size=12, italic=True, color=GRAY))]}])
    return s


# ───────────────────────── build order ─────────────────────────
def build():
    slide_title()
    slide_outline()

    divider("01", "Motivation & Problem",
            "Why fluent clinical LLMs are not yet trustworthy", TEAL)
    slide_problem()
    slide_objectives()

    divider("02", "Approach & Architecture",
            "Two trust mechanisms realised in a 12-stage pipeline", PLUM)
    slide_mechanisms()
    slide_fig("Section 02 · Pipeline", "End-to-End Inference Pipeline",
              "fig1_pipeline",
              "Twelve-stage pipeline: retrieval → graph reasoning → world-model "
              "planning → uncertainty → explanation.", PLUM)
    slide_fig("Section 02 · System", "System Architecture (Lanes A–C)",
              "fig10_system_architecture",
              "Retrieval & grounding · causal world model (10.1M params) · "
              "decision & explanation.", PLUM)
    slide_fig("Section 02 · Neural Core",
              "Detailed Neural Architecture", "fig17_architecture_detailed",
              "Context labels enter left; Shapley source-tracking produces a "
              "calibrated, attributed answer.", PLUM)
    slide_fig("Section 02 · World Model", "Latent World Model — Dataflow",
              "fig2_architecture",
              "Recurrent latent rollout z_{t+1}=f(z_t+Δz_t, a_t) enables "
              "action-conditioned planning.", PLUM)
    slide_fig("Section 02 · Capability",
              "Capability vs Existing RAG Methods", "fig11_method_capability",
              "Full / partial / no support across seven trust-relevant "
              "capabilities.", PLUM)

    divider("03", "Data & Training",
            "100,000 samples · preprocessing · convergence", SIENNA)
    slide_datasets()
    slide_fig("Section 03 · EDA", "Exploratory Data Analysis (DDXPlus)",
              "fig13_data_overview",
              "Age & sex distributions, top-15 pathologies, evidences per "
              "patient.", SIENNA)
    slide_fig("Section 03 · Cleaning", "Data Cleaning & Code Normalisation",
              "fig14_data_cleaning",
              "Record funnel (100% retained) and 504→223 evidence-code "
              "normalisation.", SIENNA)
    slide_fig("Section 03 · Preprocessing", "Preprocessing Pipeline",
              "fig18_preprocess_pipeline",
              "Clean → transform → integrate → reduce to 64-dim obs + 5-dim "
              "outcomes.", SIENNA)
    slide_fig("Section 03 · Imbalance", "Class Imbalance & Target Structure",
              "fig15_class_imbalance",
              "49 long-tailed pathologies; true-pathology rank within the "
              "differential.", SIENNA)
    slide_config()
    slide_fig("Section 03 · Convergence", "Training Convergence",
              "fig3_training_curves",
              "World-model total/reconstruction/reward losses and deep-ensemble "
              "Gaussian NLL.", SIENNA)

    divider("04", "Results & Impact",
            "Calibration, generalisation & a controlled baseline comparison",
            NAVY)
    slide_key_results()
    slide_eval_metrics()
    slide_fig("Section 04 · Calibration", "Reliability / Calibration Diagram",
              "fig4_calibration",
              "Predicted uncertainty tracks empirical error — close to the "
              "perfect-calibration line.", NAVY,
              img_box=(2.4, 1.95, 8.5, 4.5))
    slide_fig("Section 04 · Uncertainty", "Uncertainty Decomposition",
              "fig5_uncertainty",
              "Epistemic (model) vs aleatoric (data) uncertainty per case.",
              NAVY, img_box=(2.4, 1.95, 8.5, 4.5))
    slide_fig("Section 04 · Representation", "Latent Space (t-SNE)",
              "fig6_latent_tsne",
              "Evidence-Fusion-Encoder latents form clinically meaningful, "
              "pathology-aligned clusters.", NAVY,
              img_box=(2.7, 1.95, 7.9, 4.5))
    slide_fig("Section 04 · Generalisation", "Train vs Held-out Generalisation",
              "fig19_train_test",
              "World-model convergence and a train→test error gap below 0.002 "
              "MAE.", NAVY)
    slide_fig("Section 04 · Comparison",
              "Proposed Framework vs Strongest RAG Baseline",
              "fig12_method_quantitative",
              "Concept coverage, structure, evidence grounding and calibrated "
              "confidence over 10 vignettes.", NAVY)
    slide_multirag()
    slide_fig("Section 04 · Per-Vignette", "Per-Vignette Concept Coverage",
              "fig9_perprompt",
              "Concept coverage per clinical scenario (P1–P10).", NAVY)
    slide_demo()
    slide_fig("Section 04 · Significance", "Statistical Significance",
              "fig20_statistical",
              "Paired per-vignette scores and effect size (Wilcoxon p = 0.046, "
              "mean +0.27).", NAVY)

    divider("05", "Planning · Risk · Economics",
            "Engineering execution and responsible-AI considerations", GRAY)
    slide_fig("Section 05 · Economics", "Economic Analysis",
              "fig23_economic",
              "Direct project cost and marginal cost per query vs an LLM-only "
              "cloud baseline.", GRAY)
    slide_fig("Section 05 · Risk", "Risk Assessment Matrix",
              "fig22_risk_matrix",
              "Likelihood × impact for six identified project risks with "
              "mitigations.", GRAY)
    slide_fig("Section 05 · Plan", "Project Management Plan",
              "fig21_gantt",
              "Eighteen-week Gantt across literature, data, modelling, "
              "evaluation and writing.", GRAY)
    slide_impacts()

    divider("06", "Conclusion & Future Work",
            "Findings, contributions, limitations and next steps", TEAL)
    slide_contributions()
    slide_future()
    slide_conclusion()
    slide_thanks()

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Total slides: {len(prs.slides)}")


# ════════════════════════ IMPACT (societal/eth/env) ════════════════════════
def slide_impacts():
    s = new_slide()
    header(s, "Section 05 · Responsible AI", "Societal · Ethical · Environmental",
           GRAY)
    panels = [
        ("Societal Impact", TEAL, [
            "Calibrated confidence + traceable sources reduce automation bias",
            "Widens access to structured clinical reasoning",
            "Human clinician remains the final decision-maker",
        ]),
        ("Ethical Safeguards", PLUM, [
            "Only public, de-identified data (DDXPlus, MedMCQA, MedQA)",
            "Decision-support, not autonomous diagnosis",
            "Calibration + source tracking guard against confident error",
        ]),
        ("Environmental", TEAL, [
            "Compact 10.1M-param core; < 20 min training",
            "Few milliseconds of local compute per query",
            "Negligible footprint vs foundation-model training",
        ]),
    ]
    x = 0.7
    w = 3.85
    for title_, col, pts in panels:
        rect(s, x, 2.05, w, 4.5, WHITE, line=MIST, lw=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        rect(s, x, 2.05, w, 0.72, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        textbox(s, x + 0.18, 2.16, w - 0.3, 0.5,
                [{"runs": [(title_, dict(size=15.5, bold=True, color=WHITE,
                            name=HEAD_FONT))]}], anchor=MSO_ANCHOR.MIDDLE)
        bullets(s, x + 0.22, 3.0, w - 0.42, 3.4,
                [(None, p, 0) for p in pts], size=12.8, gap=12, accent=col)
        x += w + 0.2


if __name__ == "__main__":
    build()
